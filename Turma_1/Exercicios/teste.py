from pptx import Presentation
from pptx.util import Pt

# Conteúdo por década
conteudo = [
    ("Década de 1940 – As origens", [
        "Assembly (ASM): Primeiras linguagens de baixo nível usadas nos primeiros computadores.",
        "Programação com códigos binários e cartões perfurados.",
        "Início da ideia de linguagem simbólica para instruções de máquina."
    ]),
    ("Década de 1950 – Primeiras linguagens de alto nível", [
        "Fortran (1957): Focada em cálculos científicos; primeira linguagem amplamente usada.",
        "Lisp (1958): Voltada para inteligência artificial, baseada em listas.",
        "COBOL (1959): Para sistemas de gestão e negócios."
    ]),
    ("Década de 1960 – Estruturas e padronização", [
        "ALGOL (1960): Introduziu estruturas como if e while; influenciou futuras linguagens.",
        "BASIC (1964): Simples, voltada para iniciantes.",
        "Simula (1967): Primeira linguagem com orientação a objetos."
    ]),
    ("Década de 1970 – Crescimento e estruturação", [
        "C (1972): Portável e poderosa; usada em UNIX.",
        "Pascal (1970): Para ensino de boas práticas estruturadas.",
        "Prolog (1972): Linguagem lógica usada em IA."
    ]),
    ("Década de 1980 – Orientação a objetos e interfaces", [
        "C++ (1983): C com orientação a objetos.",
        "Objective-C: Combinação de C com Smalltalk (usada pela Apple).",
        "Perl (1987): Script poderoso para texto."
    ]),
    ("Década de 1990 – A internet e linguagens populares", [
        "Python (1991): Simples e legível.",
        "Java (1995): 'Write once, run anywhere'.",
        "JavaScript (1995): Para navegadores web.",
        "PHP (1995): Desenvolvimento web dinâmico."
    ]),
    ("Década de 2000 – Web, mobile e frameworks", [
        "C# (2000): Alternativa da Microsoft ao Java.",
        "Ruby on Rails (2005): Framework popular de Ruby.",
        "Swift (2014): Criada pela Apple para iOS/macOS."
    ]),
    ("Década de 2010 – Nuvem, dados e novas sintaxes", [
        "Go (2009): Desempenho e simplicidade.",
        "Rust (2010): Segurança de memória.",
        "Kotlin (2011): Alternativa moderna para Android.",
        "Popularização de linguagens para data science (R, Julia)."
    ]),
    ("Década de 2020 – IA, paralelismo e segurança", [
        "Crescimento de linguagens como Mojo, Nim, Zig.",
        "Uso intensivo de Python para IA (TensorFlow, PyTorch).",
        "Ênfase em linguagens seguras e eficientes como Rust."
    ])
]

# Criar a apresentação
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]

# Slide de título
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "História das Linguagens de Programação"
slide.placeholders[1].text = "Resumo por década"

# Slides por década
for titulo, topicos in conteudo:
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = titulo
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for topico in topicos:
        p = tf.add_paragraph()
        p.text = topico
        p.level = 0
        p.font.size = Pt(18)

# Salvar o arquivo
prs.save("Historia_Linguagens_Programacao.pptx")
print("Arquivo salvo como 'Historia_Linguagens_Programacao.pptx'")
